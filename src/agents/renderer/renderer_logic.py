from __future__ import annotations
import os
import uuid
from pathlib import Path
from typing import Dict, Any, List

from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build


from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def render_cover_slide(slide, data):
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_BLUE = RGBColor(47, 85, 151)
    

    shapes = slide.shapes

    # 상단 협회명
    box = shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
    p = box.text_frame.add_paragraph()
    p.text = "■ 부동산 마케팅 협회 ■"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.name = FONT
    p.font.color.rgb = COLOR_BLACK
    p.alignment = PP_ALIGN.CENTER

    # 메인 타이틀
    title = data.get("title", "")
    if title:
        box = shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLACK
        p.alignment = PP_ALIGN.CENTER

    # 핵심 요약 (골드)
    lead = data.get("lead", "")
    if lead:
        box = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
        p = box.text_frame.add_paragraph()
        p.text = lead
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLUE
        p.alignment = PP_ALIGN.CENTER

    # 하단 정보
    org = data.get("org", "")
    date = data.get("date", "")
    note = data.get("note", "※ 본 보고서는 내부 검토용입니다.")
    info = f" 주관: {org}\n 작성일: {date}\n\n{note}"
    box = shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(1.2))
    p = box.text_frame.add_paragraph()
    p.text = info
    p.font.size = Pt(14)
    p.font.name = FONT
    p.font.color.rgb = COLOR_BLACK
    p.alignment = PP_ALIGN.CENTER


def render_text_slide(slide, data):
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_BLUE = RGBColor(47, 85, 151)

    shapes = slide.shapes

    # --- Title ---
    title = data.get("title", "")
    if title:
        box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = 0     # 내부 여백 제거
        tf.margin_bottom = 0
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLACK

    # --- Lead (Gold) ---
    lead = data.get("lead", "")
    if lead:
        box = shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
        p = box.text_frame.add_paragraph()
        p.text = lead
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLUE

    # --- Groups ---
    groups = data.get("groups", [])
    col_w, row_h = 4.5, 2.6
    margin_x, margin_y = 0.5, 2.0

    for i, group in enumerate(groups[:4]):
        col, row = i % 2, i // 2
        x = margin_x + col * col_w
        
        # 첫 번째 줄(row == 0)만 위로 0.4인치 당겨서 리드와 간격 좁히기
        if row == 0:
            y = margin_y - 0.62
        else:
            y = (margin_y - 0.2) + row * row_h


        label = group.get("label", "")
        insight = group.get("insight", "")
        details = group.get("details", {})

        # Main box
        box = shapes.add_textbox(Inches(x), Inches(y), Inches(col_w - 0.3), Inches(row_h - 0.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_bottom = 0
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0

        # Label
        if label:
            p = tf.add_paragraph()
            p.text = label
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLACK

        # Insight (Gold)
        if insight:
            p = tf.add_paragraph()
            p.text = insight
            p.font.size = Pt(12)
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLUE

        # Details Section
        if details:
            # case 1) 최신 포맷: ["문장1", "문장2", ...]
            if isinstance(details, list):
                p = tf.add_paragraph()
                p.text = "• 세부:"
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.name = FONT
                for d in details:
                    p = tf.add_paragraph()
                    p.text = f"   - {d}"
                    p.font.size = Pt(10)
                    p.font.name = FONT

            # case 2) 예전 포맷: {"data": [...], "analysis": [...], "method": "..."}
            elif isinstance(details, dict):
                data_list = details.get("data", [])
                analysis_list = details.get("analysis", [])
                method = details.get("method")

                if data_list:
                    p = tf.add_paragraph()
                    p.text = "• 데이터:"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.name = FONT
                    for d in data_list:
                        p = tf.add_paragraph()
                        p.text = f"   - {d}"
                        p.font.size = Pt(10)
                        p.font.name = FONT

                if analysis_list:
                    p = tf.add_paragraph()
                    p.text = "• 분석:"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.name = FONT
                    for a in analysis_list:
                        p = tf.add_paragraph()
                        p.text = f"   - {a}"
                        p.font.size = Pt(10)
                        p.font.name = FONT

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def render_swot_slide(slide, data):
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_WHITE = RGBColor(255, 255, 255)
    

    title = data.get("title", "SWOT 분석")
    groups = data.get("groups", [])

    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    # Table (4 rows + header)
    rows = len(groups) + 1
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(5)).table

    # Header
    hdr_cells = table_shape.rows[0].cells
    hdr_cells[0].text = "구분"
    hdr_cells[1].text = "주요 내용"
    for cell in hdr_cells:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = FONT
            p.font.color.rgb = COLOR_WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Rows
    for i, group in enumerate(groups, start=1):
        label_cell, content_cell = table_shape.rows[i].cells
        label_cell.text = group.get("label", "")
        content_lines = group.get("details", [])
        content_cell.text = "\n".join([f"- {line}" for line in content_lines])

        for p in label_cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLACK
            p.alignment = PP_ALIGN.CENTER

        for p in content_cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLACK
            p.alignment = PP_ALIGN.LEFT

    # 테이블 스타일 간격 조정
    table_shape.columns[0].width = Inches(2.0)
    table_shape.columns[1].width = Inches(7.0)


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from utils.util import get_project_root
import uuid
def render_sliceplan_local(slice_plan: dict, output_path:str | None = None):
    if output_path is None:
        output_path = str(
            get_project_root()
            / "src" / "agents" / "renderer"
            / "temp" /f"{uuid.uuid4().hex}.pptx"
        )

    # 디렉토리가 없으면 생성
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    for slide_data in slice_plan.get("slides", []):
        layout_type = slide_data.get("layout_type", 2)
        slide = prs.slides.add_slide(blank)

        if layout_type == 1:
            render_cover_slide(slide, slide_data)
        elif layout_type == 2:
            render_text_slide(slide, slide_data)
        elif layout_type == 3:
            render_swot_slide(slide, slide_data)
        else:
            print(f"[WARN] Unknown layout_type={layout_type}")

    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    return output_path
